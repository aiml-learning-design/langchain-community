from langchain_core.document_loaders import BaseBlobParser
from pptx import Presentation  # For PPTX
from langchain_core.documents import Document
from typing import List


class PowerPointParser(BaseBlobParser):
    """
    Parser for Microsoft PowerPoint files (.pptx).

    Attempts to extract readable text from slides using `python-pptx`.
    Falls back to a basic binary read for unsupported formats like legacy `.ppt`.
    """

    def lazy_parse(self, blob) -> List[Document]:
        """
        Lazily parses a PowerPoint file and returns a list of LangChain Documents.

        Args:
            blob: A blob object representing the PowerPoint file.

        Returns:
            List[Document]: A list containing a single Document with extracted text and metadata.
        """
        if blob.path.endswith(".pptx"):
            return self._parse_pptx(blob)
        else:
            return self._parse_fallback(blob)

    def _parse_pptx(self, blob) -> List[Document]:
        """
        Parses a `.pptx` file and extracts text from all slides.

        Args:
            blob: The blob pointing to a `.pptx` file.

        Returns:
            List[Document]: A Document with concatenated slide text and metadata (including slide count).
        """
        prs = Presentation(blob.path)
        text_content = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

        return [Document(
            page_content="\n".join(text_content),
            metadata={"source": blob.path, "slides": len(prs.slides)}
        )]

    def _parse_fallback(self, blob) -> List[Document]:
        """
        Fallback parser for unsupported PowerPoint formats like `.ppt`.

        This method reads a portion of the file as raw binary text.

        Args:
            blob: A blob pointing to the file.

        Returns:
            List[Document]: A Document with raw text content as a fallback.
        """
        with open(blob.path, "rb") as f:
            raw_text = str(f.read(1024))  # Read first 1KB as fallback
        return [Document(page_content=raw_text, metadata={"source": blob.path})]